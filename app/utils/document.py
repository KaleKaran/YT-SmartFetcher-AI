import PyPDF2
from pptx import Presentation

def extract_pdf_text(file_path):
    """
    Extract text from a PDF file.
    
    Args:
        file_path (str): Path to the PDF file
        
    Returns:
        str: Extracted text
    """
    text = ""
    with open(file_path, 'rb') as file:
        pdf_reader = PyPDF2.PdfReader(file)
        for page in pdf_reader.pages:
            text += page.extract_text()
    return text

def extract_ppt_text(file_path):
    """
    Extract text from a PowerPoint presentation.
    
    Args:
        file_path (str): Path to the PPT file
        
    Returns:
        str: Extracted text
    """
    text = ""
    prs = Presentation(file_path)
    for slide in prs.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                text += shape.text + "\n"
    return text 